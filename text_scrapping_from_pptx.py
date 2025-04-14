from pptx import Presentation

# Load the presentation
presentation = Presentation(r'your_pptx_file_name.pptx')

# Open a file to save the text
with open('desire_output_file_name.txt', 'w', encoding='utf-8') as f:
    # Extract text and write it to the file
    for slide in presentation.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                f.write(shape.text + '\n')  # Write the text and add a newline

print("Text extracted and saved to 'desire_output_file_name.txt'")


##demo##
#better work on googel colab or jupyter notebook

# from pptx import Presentation

# # Load the presentation
# presentation = Presentation(r'8051_MCU.pptx')

# # Open a file to save the text with utf-8 encoding
# with open('8051_MCU.txt', 'w', encoding='utf-8') as f:
#     # Extract text and write it to the file
#     for slide in presentation.slides:
#         for shape in slide.shapes:
#             if hasattr(shape, "text"):
#                 f.write(shape.text + '\n')

# print("Text extracted and saved to '8051_MCU.txt'")


